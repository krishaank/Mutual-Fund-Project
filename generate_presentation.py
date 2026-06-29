"""
Generates the 12-slide Capstone Presentation using python-pptx.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Slide layouts
    title_slide_layout = prs.slide_layouts[0]
    title_content_layout = prs.slide_layouts[1]
    
    # Custom Theme Colors
    BLUE = RGBColor(31, 58, 147)
    
    def set_title(slide, text):
        title = slide.shapes.title
        title.text = text
        title.text_frame.paragraphs[0].font.color.rgb = BLUE
        title.text_frame.paragraphs[0].font.bold = True
    
    # 1. Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Bluestock Mutual Fund Analytics"
    title.text_frame.paragraphs[0].font.color.rgb = BLUE
    subtitle.text = "Capstone Project Presentation\\nData Engineering, Advanced Analytics & BI"
    
    # 2. Problem & Objective
    slide = prs.slides.add_slide(title_content_layout)
    set_title(slide, "Problem & Objective")
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Problem Statement:"
    p = tf.add_paragraph()
    p.text = "- Lack of consolidated insights across 40 Mutual Fund schemes."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "- Inability to decouple market returns (Beta) from manager skill (Alpha)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Objectives:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "- Design an automated ETL pipeline and Star Schema data warehouse."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "- Compute advanced risk-adjusted metrics (Sharpe, CVaR, Max Drawdown)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "- Build an interactive BI Dashboard for executive decision-making."
    p.level = 1

    # 3. Data Sources & Architecture
    slide = prs.slides.add_slide(title_content_layout)
    set_title(slide, "Data Architecture")
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Data Sources:"
    p = tf.add_paragraph()
    p.text = "- 10 Raw CSV datasets spanning NAV, Transactions, and Portfolio Holdings."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "ETL & Star Schema:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "- Powered by Python (Pandas) and SQLite."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "- Fact Tables: fact_nav, fact_transactions, fact_aum, fact_performance."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "- Dimension Tables: dim_fund, dim_date."
    p.level = 1

    # 4. EDA Highlights 1
    slide = prs.slides.add_slide(title_content_layout)
    set_title(slide, "EDA Highlights: Macro Growth")
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Total Assets Under Management (AUM):"
    p = tf.add_paragraph()
    p.text = "- The industry hit a historic milestone of ₹81 Lakh Crores."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Retail Participation Surge:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "- Systematic Investment Plan (SIP) inflows peaked at ₹31,002 Crores per month."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "- Total investor folios expanded to 26.12 Crores."
    p.level = 1

    # 5. EDA Highlights 2
    slide = prs.slides.add_slide(title_content_layout)
    set_title(slide, "EDA Highlights: Demographics")
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Geographical Distribution:"
    p = tf.add_paragraph()
    p.text = "- Maharashtra and Gujarat lead in total SIP contributions."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Age Demographics:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "- 25-35 year olds constitute the largest volume of SIP transactions."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "- 45-55 year olds maintain the highest average ticket size per SIP."
    p.level = 1

    # 6. Performance Metrics 1
    slide = prs.slides.add_slide(title_content_layout)
    set_title(slide, "Performance Analytics: Returns")
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Compound Annual Growth Rate (CAGR):"
    p = tf.add_paragraph()
    p.text = "- Calculated 1-Year, 3-Year, and 5-Year CAGR for all 40 schemes."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Alpha & Beta Extraction:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "- Utilized OLS Regression via scipy.stats against the Nifty 100 benchmark."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "- Quantified true manager skill (Alpha) isolated from market correlation (Beta)."
    p.level = 1

    # 7. Performance Metrics 2
    slide = prs.slides.add_slide(title_content_layout)
    set_title(slide, "Performance Analytics: Risk")
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Risk-Adjusted Ratios:"
    p = tf.add_paragraph()
    p.text = "- Sharpe Ratio (utilizing a 6.5% risk-free rate) & Sortino Ratio."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Tail Risk Analysis:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "- Historical 95% Value at Risk (VaR) & Conditional VaR (CVaR)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "- Calculated Herfindahl-Hirschman Index (HHI) identifying hidden sector concentration."
    p.level = 1

    # 8. Dashboard 1
    slide = prs.slides.add_slide(title_content_layout)
    set_title(slide, "BI Dashboard: Industry Overview")
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Built natively in Python (Streamlit):"
    p = tf.add_paragraph()
    p.text = "- Direct SQLite ODBC connection ensures real-time queries."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "- Interactive KPI cards tracking total industry assets and flows."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "- Visual breakdown of AMC market share and total AUM trends."
    p.level = 1

    # 9. Dashboard 2
    slide = prs.slides.add_slide(title_content_layout)
    set_title(slide, "BI Dashboard: Analytics & Cohorts")
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Fund Performance Matrix:"
    p = tf.add_paragraph()
    p.text = "- Dynamic Risk vs Return scatter plots (Bubble size = Expense Ratio)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Investor Analytics:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "- Behavioral cohort tracking (Lifetime Value by starting year)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "- Cross-filtering enabled across demographics, category, and fund house."
    p.level = 1

    # 10. Key Findings
    slide = prs.slides.add_slide(title_content_layout)
    set_title(slide, "Key Insights & Findings")
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Sticky Capital Hypothesis Validated:"
    p = tf.add_paragraph()
    p.text = "- SIP inflows remained highly resilient despite the 2024 Nifty 50 corrections."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Manager Underperformance:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "- Rolling 90-Day Sharpe ratios revealed significant manager underperformance during volatile quarters."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "The Composite Scorecard:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "- Algorithmic 0-100 grading successfully penalized funds with high expense ratios and excessive max drawdowns."
    p.level = 1

    # 11. Limitations & Recommendations
    slide = prs.slides.add_slide(title_content_layout)
    set_title(slide, "Limitations & Recommendations")
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Limitations:"
    p = tf.add_paragraph()
    p.text = "- 3-year history is insufficient for full economic cycle stress testing."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "- Lack of individual stock transaction data restricts deepest attribution analysis."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Recommendations:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "- Implement the Churn Prediction algorithm to proactively target 'at-risk' SIP investors (gaps > 35 days)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "- Deploy the Fund Recommender Engine onto the retail mobile application."
    p.level = 1

    # 12. Thank You
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Thank You"
    title.text_frame.paragraphs[0].font.color.rgb = BLUE
    subtitle.text = "Open for Questions & Feedback"
    
    prs.save('Bluestock_MF_Presentation.pptx')
    print("Bluestock_MF_Presentation.pptx generated successfully.")

if __name__ == '__main__':
    create_presentation()
